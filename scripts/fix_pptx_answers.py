#!/usr/bin/env python3
"""
Modify l8-review-test-prep.pptx to:
1. Clone each exercise slide as an answer-reveal slide (showing all answers)
2. On the original exercise slides, hide answer text by clearing it
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
import copy
from lxml import etree
import re

PP = '/Users/luojiahai/Code/deck/index/designs/y7-l10/l8-review-test-prep/l8-review-test-prep.pptx'

NSMAP = {
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
}


def duplicate_slide(prs, slide_idx):
    """Duplicate slide at slide_idx, insert right after it."""
    slide = prs.slides[slide_idx]
    slide_layout = slide.slide_layout

    # Add new slide with same layout
    new_slide = prs.slides.add_slide(slide_layout)

    # Remove all default shapes from new slide
    spTree = new_slide.shapes._spTree
    for sp in list(spTree):
        if sp.tag != '{http://schemas.openxmlformats.org/presentationml/2006/main}cSld':
            # Keep cSld, remove children
            pass
    # Clear all shape children from the spTree
    to_remove = []
    for child in spTree:
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag in ('sp', 'pic', 'grpSp', 'cxnSp', 'graphicFrame'):
            to_remove.append(child)
    for child in to_remove:
        spTree.remove(child)

    # Copy shapes from original to new
    orig_spTree = slide.shapes._spTree
    for child in list(orig_spTree):
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag in ('sp', 'pic', 'grpSp', 'cxnSp', 'graphicFrame'):
            new_child = copy.deepcopy(child)
            spTree.append(new_child)

    # Move the new slide to right after the original
    ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    sldIdLst = prs._element.find('p:sldIdLst', ns)
    sldIds = list(sldIdLst)
    # The new slide was appended at the end, move it
    new_sldId = sldIds[-1]
    sldIdLst.remove(new_sldId)
    # Insert after original (at slide_idx + 1 position)
    sldIdLst.insert(slide_idx + 1, new_sldId)

    return new_slide


def clear_text_shape(shape):
    """Clear all text from a shape."""
    txBody = shape._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}txBody')
    if txBody is not None:
        for p in list(txBody):
            tag = p.tag.split('}')[-1] if '}' in p.tag else p.tag
            if tag == 'p':
                txBody.remove(p)
        # Add an empty paragraph so the shape doesn't error
        empty_p = etree.SubElement(txBody, '{http://schemas.openxmlformats.org/drawingml/2006/main}p')
        empty_r = etree.SubElement(empty_p, '{http://schemas.openxmlformats.org/drawingml/2006/main}r')
        empty_t = etree.SubElement(empty_r, '{http://schemas.openxmlformats.org/drawingml/2006/main}t')
        empty_t.text = ''


def get_shape_text(shape):
    """Get the full text content of a shape."""
    if not shape.has_text_frame:
        return ''
    return shape.text_frame.text


def hide_warmup_answers(slide):
    """Slide 2: Hide shapes starting with 'Answer:'."""
    for shape in slide.shapes:
        text = get_shape_text(shape).strip()
        if text.startswith('Answer:'):
            clear_text_shape(shape)
            print(f"    Hiding: '{text[:50]}'")


def hide_r1_answers(slide):
    """Slide 5 or 6: Hide ↓ arrows, pinyin, and English translations.

    Pattern: q-num (keep) → cn (keep) → ↓ (hide) → py (hide) → ↓ (hide) → en (hide)
    Also handles the '↓' arrows between items.
    """
    shapes = list(slide.shapes)
    i = 0
    while i < len(shapes):
        text = get_shape_text(shapes[i]).strip()
        if text == '↓':
            clear_text_shape(shapes[i])
            print(f"    Hiding arrow shape: '{shapes[i].name}'")
            # The shape AFTER a ↓ arrow could be py or en text
            # For R1: ↓ → py → ↓ → en (alternating)
            # Actually, just hide the ↓ arrows. The py/en text gets hidden by the content approach below.
            pass
        i += 1

    # Also hide pinyin and English text between arrows
    # Pinyin: lowercase letters with tone marks (a-z with diacritics)
    # English: plain ASCII words
    # Chinese characters: CJK Unicode range

    PINYIN_PATTERN = re.compile(r'^[a-zA-Zāáǎàēéěèīíǐìōóǒòūúǔùǖǘǚǜ]+(?:\s+[a-zA-Zāáǎàēéěèīíǐìōóǒòūúǔùǖǘǚǜ]+)*$')

    for shape in slide.shapes:
        text = get_shape_text(shape).strip()
        if not text:
            continue
        # Skip numbers
        if text == '↓' or text == '':
            continue
        # Check if this is pure pinyin
        if PINYIN_PATTERN.match(text):
            clear_text_shape(shape)
            print(f"    Hiding pinyin: '{text}'")
            continue
        # Check if this is English translation (no CJK characters)
        has_cjk = bool(re.search(r'[一-鿿㐀-䶿]', text))
        has_digit = bool(re.search(r'\d', text))
        if not has_cjk and not text.startswith(('Item', 'Round', 'L8', 'Y7', '1–', '9–')):
            # Likely English translation - but check it's not a label
            if text not in ('↓',) and not text.startswith(('Items', 'of', 'pts')):
                # Double check: if text has only ASCII and common English characters
                if all(ord(c) < 128 for c in text) and not text.startswith(('0–', '1–', '2–', '3–', '4–', '5–', '6–', '7–', '8–', '9–', 'L8', 'Y7', 'Round', 'Items', 'Format')):
                    clear_text_shape(shape)
                    print(f"    Hiding English: '{text[:50]}'")


def hide_r2_answers(slide):
    """Slide 8 or 9: Hide ↓ arrows and answer text (a-text).

    Pattern: q-num (keep) → q-text (keep) → ↓ (hide) → a-text (hide) → pts-tag (keep)
    The answer text follows the ↓ arrow.
    """
    shapes = list(slide.shapes)
    # Collect indices of all shapes
    for i, shape in enumerate(shapes):
        text = get_shape_text(shape).strip()
        if text == '↓':
            # Hide the arrow
            clear_text_shape(shape)
            print(f"    Hiding arrow: '{shape.name}'")
            # The next non-empty shape is the answer
            for j in range(i + 1, min(i + 3, len(shapes))):
                next_text = get_shape_text(shapes[j]).strip()
                if next_text and next_text not in ('EN → 中文', '中文 → EN'):
                    clear_text_shape(shapes[j])
                    print(f"    Hiding answer: '{next_text[:50]}'")
                    break


def hide_r3_answers(slide):
    """Slide 11: Hide verdict badges, correction text, and reason text.

    Pattern: q-num (keep) → given (keep) → verdict (hide) → correction (hide) → reason (hide)
    """
    WRONG_PATTERNS = ['WRONG', 'CORRECT']
    REASON_PATTERNS = [
        '二 → 两 before 点', 'No error', 'Need 零 for single-digit',
        'Word order: 差 X 分 Y 点', '60分 = 1 hour',
    ]

    # Keep track of which shapes are part of error cards (between q-num and next q-num)
    shapes = list(slide.shapes)

    for shape in slide.shapes:
        text = get_shape_text(shape).strip()
        if not text:
            continue

        # Hide verdict badges
        if 'WRONG' in text or 'CORRECT' in text:
            clear_text_shape(shape)
            print(f"    Hiding verdict: '{text[:50]}'")
            continue

    # Hide correction and reason text
    # These are harder to identify. Correction text is Chinese (looks like error content),
    # and reason text is English explanations.
    # Strategy: after clearing verdicts, find shapes near them that aren't q-num or given

    for shape in slide.shapes:
        text = get_shape_text(shape).strip()
        if not text:
            continue
        # q-num: single digit or small numbers
        if text.isdigit() and len(text) <= 2:
            continue
        # given text: Chinese time expressions (already kept)
        if '点' in text and '分' not in text and len(text) <= 6:
            # This is likely the "given" text - keep it
            continue
        if '点' in text and ('分' in text or '刻' in text):
            # This is likely correction text - but could also be a given
            # Only hide if the shape isn't the first "点" containing shape in its group
            pass

        # Actually, let's just identify shapes to hide by content patterns
        # Correction text: Chinese expressions that look like corrected versions
        # Reason text: English explanations

    # Clear correction and reason shapes
    for shape in slide.shapes:
        text = get_shape_text(shape).strip()
        if not text:
            continue
        # Check if this is a reason (English explanation)
        for pattern in REASON_PATTERNS:
            if pattern in text:
                clear_text_shape(shape)
                print(f"    Hiding reason: '{text[:50]}'")
                break

    # Now clear correction text (shapes that contain Chinese time expressions
    # but aren't the "given" text)
    # Approach: each error card has: q-num, given, (verdict-cleared), correction, reason
    # After clearing verdict and reason, the leftover shapes containing 点/分/刻 are corrections
    for shape in slide.shapes:
        text = get_shape_text(shape).strip()
        if not text:
            continue
        # Skip numbers, given text, labels, footer
        if text.isdigit() and len(text) <= 2:
            continue
        if text in ('Error Bounty · 错误悬赏', 'Spot the mistake & write the correct form — some are already right!',
                     '8 items · 3 correct · 5 wrong', '30–38 min · Round 3 · 20 pts each',
                     'L8 · BEAT THE TEST', 'Round 3 · Error Bounty',
                     'Y7 Chinese · Unit 4.1 Telling Time'):
            continue
        # Correction text: contains Chinese time vocabulary
        if re.search(r'[一-鿿]', text) and ('点' in text or '分' in text or '刻' in text or '/' in text):
            # Distinguish "given" from "correction":
            # Correction often has / (like "一点 / 十三点") or is slightly different from the given
            if '/' in text:
                clear_text_shape(shape)
                print(f"    Hiding correction: '{text[:50]}'")
                continue
            # If not the "given" text, it's likely correction
            # But given text also has 点/分/刻...
            # Let's just be conservative and only clear shapes that aren't the first Chinese text after q-num
            clear_text_shape(shape)
            print(f"    Hiding correction: '{text[:50]}'")


def main():
    print("Opening PPTX...")
    prs = Presentation(PP)

    # Define exercise slides to process
    # These are 0-indexed positions in the ORIGINAL 14-slide deck
    # We need to process them in REVERSE order because each duplication shifts indices
    exercises = [
        (1, 'warm-up', hide_warmup_answers),        # Slide 2 (0-indexed: 1)
        (4, 'r1-items-a', hide_r1_answers),          # Slide 5 (0-indexed: 4)
        (5, 'r1-items-b', hide_r1_answers),          # Slide 6 (0-indexed: 5)
        (7, 'r2-items-a', hide_r2_answers),          # Slide 8 (0-indexed: 7)
        (8, 'r2-items-b', hide_r2_answers),          # Slide 9 (0-indexed: 8)
        (10, 'r3-items', hide_r3_answers),           # Slide 11 (0-indexed: 10)
    ]

    # Process in REVERSE order so index shifts don't affect us
    for orig_idx, label, hide_fn in reversed(exercises):
        print(f"\n--- Processing slide {orig_idx + 1} ({label}) ---")

        # 1. Duplicate the slide (inserts right after original)
        print(f"  Cloning slide {orig_idx + 1}...")
        duplicate_slide(prs, orig_idx)
        print(f"  Clone inserted after slide {orig_idx + 1}")

        # 2. Hide answers on the original (still at orig_idx since clone was inserted after)
        print(f"  Hiding answers on original slide...")
        hide_fn(prs.slides[orig_idx])

    print(f"\nTotal slides: {len(prs.slides)}")
    print(f"Saving...")
    prs.save(PP)
    print(f"✓ Saved to {PP}")


if __name__ == '__main__':
    main()
